from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Presentation setup ───────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ── Colour palette (matches powerpoint_theme.md) ─────────────────────
BG      = RGBColor(0xF7, 0xF9, 0xFB)  # off-white background
ACCENT  = RGBColor(0x4A, 0x90, 0xE2)  # muted blue (primary accent)
TEAL    = RGBColor(0x5F, 0xA8, 0xA0)  # soft teal (secondary accent)
AMBER   = RGBColor(0xD4, 0x8F, 0x0A)  # warm amber (highlight)
TEXT    = RGBColor(0x1A, 0x1A, 0x1A)  # near-black
DIM     = RGBColor(0x6B, 0x72, 0x80)  # muted gray
STEP_BG = RGBColor(0xED, 0xF2, 0xF7)  # light blue-gray fill
DIVIDER = RGBColor(0xD1, 0xD9, 0xE3)  # subtle separator
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
FONT    = "Segoe UI"

# ── Helpers ──────────────────────────────────────────────────────────
def set_bg(slide):
    fill = slide.background.fill
    fill.solid(); fill.fore_color.rgb = BG

def tb(slide, text, x, y, w, h, sz=13, bold=False, italic=False,
       color=None, align=PP_ALIGN.LEFT, font=FONT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = box.text_frame; tf.word_wrap = True
    p   = tf.paragraphs[0]; p.alignment = align
    r   = p.add_run()
    r.text = text; r.font.size = Pt(sz); r.font.name = font
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color or TEXT
    return box

def solid_rect(slide, x, y, w, h, fill, border=None, bw=Pt(0.75)):
    shp = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if border: shp.line.color.rgb = border; shp.line.width = bw
    else: shp.line.fill.background()
    return shp

def round_rect(slide, x, y, w, h, fill, border=None, bw=Pt(0.75)):
    shp = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if border: shp.line.color.rgb = border; shp.line.width = bw
    else: shp.line.fill.background()
    return shp

def hdiv(slide, x, y, w, color=None):
    solid_rect(slide, x, y, w, 0.018, fill=color or DIVIDER)

def slide_header(slide, label):
    """Top accent bar + small label used on every content slide."""
    solid_rect(slide, 0, 0, 13.33, 0.05, ACCENT)
    tb(slide, label, 0.5, 0.12, 4.0, 0.30, sz=9, bold=True, color=DIM)
    hdiv(slide, 0.5, 0.48, 12.33)

def slide_footer(slide):
    solid_rect(slide, 0, 7.45, 13.33, 0.05, ACCENT)


# ════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); set_bg(sl)
solid_rect(sl, 0, 0, 13.33, 0.05, ACCENT)
solid_rect(sl, 0, 2.8, 13.33, 2.2, fill=STEP_BG)

tb(sl, "BIRCH WILSON", 0, 0.9, 13.33, 1.6,
   sz=72, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
tb(sl, "Sprint 1 Planning  ·  Sprint Backlog & Definition of Done",
   1, 3.0, 11.33, 0.6, sz=20, color=TEXT, align=PP_ALIGN.CENTER)
tb(sl, "MET CS 634 — Agile Software Development",
   1, 3.68, 11.33, 0.44, sz=13, color=DIM, align=PP_ALIGN.CENTER)
tb(sl, "John LaCroix  ·  April 2026",
   1, 4.18, 11.33, 0.44, sz=12, color=DIM, align=PP_ALIGN.CENTER)
slide_footer(sl)


# ════════════════════════════════════════════════════════════════════
# SLIDE 2 — SPRINT GOAL
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); set_bg(sl)
slide_header(sl, "SPRINT GOAL")

# Sprint number pill
round_rect(sl, 10.8, 0.10, 2.3, 0.30, fill=ACCENT)
tb(sl, "SPRINT 1", 10.8, 0.10, 2.3, 0.30,
   sz=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Headline insight
tb(sl, "Give Dan a Story to Tell. Give Jake a Reason to Call.",
   0.5, 0.60, 12.33, 0.72, sz=24, bold=True, color=TEXT)

# Goal statement box
round_rect(sl, 0.5, 1.42, 12.33, 1.08, fill=STEP_BG, border=ACCENT)
tb(sl, "By the end of Sprint 1, Dan (underground blogger) can discover, research, and sign up — "
       "and Jake (indie producer) can evaluate the band's sound, verify they're actively gigging, "
       "and submit a recording inquiry. Both personas complete their journey end-to-end on a live site.",
   0.70, 1.52, 11.93, 0.88, sz=12, color=TEXT)

hdiv(sl, 0.5, 2.66, 12.33)

# Two persona columns
# ── Dan column (left, TEAL)
solid_rect(sl, 0.5, 2.80, 5.9, 0.46, fill=TEAL)
tb(sl, "DAN  ·  The Hardcore Blogger", 0.5, 2.80, 5.9, 0.46,
   sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
round_rect(sl, 0.5, 3.28, 5.9, 0.36, fill=RGBColor(0xE8, 0xF4, 0xF3), border=TEAL, bw=Pt(0.75))
tb(sl, "Age 30  ·  Allston, MA  ·  Freelance Underground Music Blogger",
   0.68, 3.34, 5.54, 0.24, sz=8, italic=True, color=TEAL)
dan_goals = [
    "Find the site via search ('underground rock bands MA')",
    "Read band bio & origin story on the About page",
    "Stream Spotify tracks to assess sound and authenticity",
    "Sign up for show alerts so he can cover them live",
]
for i, g in enumerate(dan_goals):
    round_rect(sl, 0.5, 3.72 + i * 0.56, 5.9, 0.50, fill=RGBColor(0xF5, 0xF8, 0xFC), border=TEAL, bw=Pt(0.5))
    tb(sl, f"→  {g}", 0.68, 3.80 + i * 0.56, 5.54, 0.34, sz=9, color=TEXT)

# ── Jake column (right, AMBER)
solid_rect(sl, 6.93, 2.80, 5.9, 0.46, fill=AMBER)
tb(sl, "JAKE  ·  The Indie Producer", 6.93, 2.80, 5.9, 0.46,
   sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
round_rect(sl, 6.93, 3.28, 5.9, 0.36, fill=RGBColor(0xFD, 0xF5, 0xE0), border=AMBER, bw=Pt(0.75))
tb(sl, "West Springfield, MA  ·  Ghost Hit Recording  ·  Indie & Rock Specialist",
   7.11, 3.34, 5.54, 0.24, sz=8, italic=True, color=AMBER)
jake_goals = [
    "Listen to Spotify embeds; assess production quality & fit",
    "Read About page — verify band aesthetic matches his niche",
    "Check Shows page to confirm the band is actively gigging",
    "Submit a recording collaboration inquiry via contact form",
]
for i, g in enumerate(jake_goals):
    round_rect(sl, 6.93, 3.72 + i * 0.56, 5.9, 0.50, fill=RGBColor(0xF5, 0xF8, 0xFC), border=AMBER, bw=Pt(0.5))
    tb(sl, f"→  {g}", 7.11, 3.80 + i * 0.56, 5.54, 0.34, sz=9, color=TEXT)

# Divider column center gap
solid_rect(sl, 6.44, 2.80, 0.46, 4.00, fill=BG)
tb(sl, "vs", 6.44, 4.60, 0.46, 0.36, sz=9, bold=True, color=DIM, align=PP_ALIGN.CENTER)

# Scope line
hdiv(sl, 0.5, 6.68, 12.33)
tb(sl, "Scope:  17 user stories  ·  7 epics  ·  T-shirt sized (XS → L)  ·  Release 1 — Core Foundation",
   0.5, 6.76, 12.33, 0.36, sz=10, color=DIM)

slide_footer(sl)


# ════════════════════════════════════════════════════════════════════
# SLIDE 3 — TRELLO BOARD
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); set_bg(sl)
slide_header(sl, "SPRINT BACKLOG")

tb(sl, "Sprint 1 Backlog lives in Trello",
   0.5, 0.60, 12.33, 0.72, sz=28, bold=True, color=TEXT)

# Board link card
round_rect(sl, 0.5, 1.46, 12.33, 0.84, fill=STEP_BG, border=ACCENT)
tb(sl, "Trello Board:", 0.72, 1.56, 2.40, 0.30, sz=11, bold=True, color=DIM)
tb(sl, "https://trello.com/b/M4SRMaOV/cs634lacroixjohn",
   3.00, 1.56, 9.60, 0.30, sz=11, bold=True, color=ACCENT)
tb(sl, "CS634LaCroixJohn  ·  Sprint 1 Backlog list",
   0.72, 1.88, 11.50, 0.28, sz=9, italic=True, color=DIM)

hdiv(sl, 0.5, 2.50, 12.33)

# Board structure breakdown (two columns)
tb(sl, "Board Structure", 0.5, 2.62, 5.8, 0.36, sz=13, bold=True, color=TEXT)
tb(sl, "Sprint 1 Contents", 7.0, 2.62, 5.8, 0.36, sz=13, bold=True, color=TEXT)

lists_left = [
    ("Trello Starter Guide",          "Orientation & onboarding notes"),
    ("Product Backlog",               "Full prioritized product backlog"),
    ("Sprint 1 Backlog",              "Active sprint — 17 user stories"),
    ("Release 1 — Core Foundation",   "Release 1 epics & backlog items"),
    ("Release 2 — Fan Engagement",    "Post-MVP release items"),
    ("Release 3",                     "Future scope"),
]
for i, (name, desc) in enumerate(lists_left):
    y = 3.08 + i * 0.62
    is_sprint = name.startswith("Sprint 1")
    fill  = STEP_BG if is_sprint else RGBColor(0xF5, 0xF8, 0xFC)
    bdr   = ACCENT  if is_sprint else DIVIDER
    bw    = Pt(1.25) if is_sprint else Pt(0.75)
    round_rect(sl, 0.5, y, 5.8, 0.54, fill=fill, border=bdr, bw=bw)
    tb(sl, name, 0.70, y + 0.04, 3.60, 0.22, sz=9,
       bold=is_sprint, color=ACCENT if is_sprint else TEXT)
    tb(sl, desc, 0.70, y + 0.28, 5.40, 0.22, sz=8, color=DIM)

stats = [
    ("17",  "User Stories (US-01 → US-17)"),
    ("7",   "Epics  (Homepage, Bio, Music, Instagram, Contact, Shows, Email)"),
    ("2",   "Anchor cards  (Sprint Goal  +  Definition of Done)"),
    ("4",   "T-shirt sizes used  (XS · S · M · L)"),
    ("5",   "Labels configured  (XS · S · M · L · XL)"),
]
for i, (num, label) in enumerate(stats):
    y = 3.08 + i * 0.72
    round_rect(sl, 7.0, y, 5.8, 0.62, fill=RGBColor(0xF5, 0xF8, 0xFC), border=DIVIDER)
    tb(sl, num,   7.18, y + 0.06, 0.72, 0.50, sz=22, bold=True, color=ACCENT)
    tb(sl, label, 7.90, y + 0.14, 4.70, 0.36, sz=9,  color=TEXT)

slide_footer(sl)


# ════════════════════════════════════════════════════════════════════
# SLIDE 4 — DEFINITION OF DONE
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); set_bg(sl)
slide_header(sl, "DEFINITION OF DONE")

round_rect(sl, 10.8, 0.10, 2.3, 0.30, fill=TEAL)
tb(sl, "SPRINT 1", 10.8, 0.10, 2.3, 0.30,
   sz=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

tb(sl, "A Story Is Done When All Five Gates Are Cleared",
   0.5, 0.60, 12.33, 0.72, sz=26, bold=True, color=TEXT)

dod_items = [
    ("01", "Acceptance Criteria",
     "All acceptance criteria on the Trello card checklist are checked off by the developer "
     "and confirmed by the product owner before the card moves to Done."),
    ("02", "Deployed & Functional",
     "The feature is visible and working on the deployed site. No broken links, "
     "missing images, or JavaScript console errors are present."),
    ("03", "Cross-Device Tested",
     "The feature has been tested on both desktop (1280px) and mobile (375px) viewports. "
     "Layout, readability, and interactions are verified on both."),
    ("04", "No Critical Defects",
     "No P1 or P2 bugs remain open against the story. Minor cosmetic issues may be "
     "logged as follow-on cards but do not block Done status."),
    ("05", "Card Closed in Sprint",
     "The Trello card has been moved to the Done column before the sprint ends. "
     "Any incomplete work is split into a new card and returned to the backlog."),
]

for i, (num, title, body) in enumerate(dod_items):
    y = 1.52 + i * 1.12
    # Number badge
    round_rect(sl, 0.5, y, 0.58, 0.96, fill=TEAL)
    tb(sl, num, 0.5, y + 0.22, 0.58, 0.52,
       sz=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Content box
    round_rect(sl, 1.18, y, 11.65, 0.96, fill=STEP_BG, border=DIVIDER)
    tb(sl, title, 1.38, y + 0.06, 11.20, 0.28, sz=11, bold=True, color=TEXT)
    tb(sl, body,  1.38, y + 0.36, 11.20, 0.54, sz=9,  color=DIM)

slide_footer(sl)


# ════════════════════════════════════════════════════════════════════
# SLIDE 5 — REFERENCES (APA 7th Edition)
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); set_bg(sl)
slide_header(sl, "REFERENCES")

tb(sl, "References",
   0.5, 0.60, 12.33, 0.64, sz=26, bold=True, color=TEXT)

hdiv(sl, 0.5, 1.30, 12.33, ACCENT)

# Each reference: (hanging-indent line 1, continuation line 2, URL line 3)
refs = [
    (
        "Cohn, M. (2023, June 20). The sprint goal: What it is and how it can help.",
        "    Mountain Goat Software.",
        "    https://www.mountaingoatsoftware.com/blog/the-sprint-goal-what-it-is-and-how-it-can-help",
    ),
    (
        "Easy Agile Team. (2024, July 2). How to write user stories in agile software development.",
        "    Easy Agile.",
        "    https://www.easyagile.com/blog/how-to-write-good-user-stories-in-agile-software-development",
    ),
    (
        "Kramer, N. (2024, August 11). T-shirt sizing in agile: Guide 2024.",
        "    daily.dev.",
        "    https://daily.dev/blog/t-shirt-sizing-in-agile-guide-2024",
    ),
]

for i, (line1, line2, line3) in enumerate(refs):
    y = 1.50 + i * 1.52
    round_rect(sl, 0.5, y, 12.33, 1.36, fill=STEP_BG, border=DIVIDER)
    tb(sl, line1, 0.70, y + 0.10, 11.93, 0.34, sz=11, color=TEXT)
    tb(sl, line2, 0.70, y + 0.46, 11.93, 0.28, sz=11, color=TEXT)
    tb(sl, line3, 0.70, y + 0.78, 11.93, 0.46, sz=10, italic=True, color=ACCENT)

# APA note
tb(sl, "All references formatted per APA 7th Edition.",
   0.5, 6.90, 12.33, 0.34, sz=9, italic=True, color=DIM)

slide_footer(sl)


# ── Save ─────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "LacroixJohn_Assignment4_Sprint1_BirchWilson.pptx")
prs.save(out)
print(f"Saved: {out}")
