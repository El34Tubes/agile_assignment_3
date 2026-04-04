from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Presentation setup ───────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ── Colour palette (product-focused light theme) ─────────────────────
BG            = RGBColor(0xF7, 0xF9, 0xFB)  # off-white background
ACCENT        = RGBColor(0x4A, 0x90, 0xE2)  # muted blue (primary accent)
TEXT          = RGBColor(0x1A, 0x1A, 0x1A)  # near-black primary text
DIM           = RGBColor(0x6B, 0x72, 0x80)  # muted gray secondary text
STEP_BG       = RGBColor(0xED, 0xF2, 0xF7)  # light blue-gray
ACT_BG        = RGBColor(0xF5, 0xF8, 0xFC)  # very light fill
BAR_BG        = RGBColor(0xEE, 0xF2, 0xF6)  # bar background
DAN_C         = RGBColor(0x5F, 0xA8, 0xA0)  # soft teal  (Dan)
JAKE_C        = RGBColor(0xD4, 0x8F, 0x0A)  # warm amber (Jake)
DAN_DARK      = RGBColor(0xE8, 0xF4, 0xF3)  # light teal bg
JAKE_DARK     = RGBColor(0xFD, 0xF5, 0xE0)  # light amber bg
DIVIDER       = RGBColor(0xD1, 0xD9, 0xE3)  # subtle separator
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)  # pure white (pill text)
FONT          = "Segoe UI"

# ── Flow slide layout constants ──────────────────────────────────────
SB_W   = 1.88
CX     = 1.96
CW     = 11.30
COL_GAP = 0.10
COL_W  = (CW - 3 * COL_GAP) / 4
COL_XS = [CX + i * (COL_W + COL_GAP) for i in range(4)]

CAP_Y  = 0.52;  CAP_H  = 0.65
STEP_Y = 1.24;  STEP_H = 0.60
ACT_Y  = 1.90;  ACT_H  = 0.46
ACT_GAP = 0.05
EN_Y   = 4.42;  EN_H   = 0.70
NF_Y   = 5.18;  NF_H   = 0.62
DT_Y   = 5.86;  DT_H   = 0.70

# ── Helpers ──────────────────────────────────────────────────────────
def set_bg(slide, color=None):
    fill = slide.background.fill
    fill.solid(); fill.fore_color.rgb = color or BG

def tb(slide, text, x, y, w, h, sz=13, bold=False, italic=False,
       color=None, align=PP_ALIGN.LEFT, wrap=True, font=FONT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = box.text_frame; tf.word_wrap = wrap
    p   = tf.paragraphs[0]; p.alignment = align
    r   = p.add_run()
    r.text = text; r.font.size = Pt(sz); r.font.name = font
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color or TEXT
    return box

def solid_rect(slide, x, y, w, h, fill, border=None, bw=Pt(1.0)):
    shp = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if border: shp.line.color.rgb = border; shp.line.width = bw
    else: shp.line.fill.background()
    return shp

def round_rect(slide, x, y, w, h, fill, border=None, bw=Pt(1.0)):
    shp = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if border: shp.line.color.rgb = border; shp.line.width = bw
    else: shp.line.fill.background()
    return shp

def shape_with_text(slide, shape_type, x, y, w, h, fill, border,
                    label_line1, label_line2, sz1=7, sz2=10,
                    c1=None, c2=None, bw=Pt(0.75), ml=0.07):
    shp = slide.shapes.add_shape(shape_type,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border; shp.line.width = bw
    tf = shp.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(ml); tf.margin_right  = Inches(0.07)
    tf.margin_top  = Inches(0.05); tf.margin_bottom = Inches(0.03)
    if label_line1:
        p1 = tf.paragraphs[0]; r1 = p1.add_run()
        r1.text = label_line1; r1.font.size = Pt(sz1); r1.font.name = FONT
        r1.font.bold = True; r1.font.color.rgb = c1 or ACCENT
        if label_line2:
            p2 = tf.add_paragraph(); r2 = p2.add_run()
            r2.text = label_line2; r2.font.size = Pt(sz2); r2.font.name = FONT
            r2.font.color.rgb = c2 or TEXT
    elif label_line2:
        p1 = tf.paragraphs[0]; r1 = p1.add_run()
        r1.text = label_line2; r1.font.size = Pt(sz2); r1.font.name = FONT
        r1.font.color.rgb = c2 or TEXT
    return shp

def hdiv(slide, x, y, w, color=None):
    solid_rect(slide, x, y, w, 0.018, fill=color or DIVIDER)

# ── Sidebar helper for flow slides ───────────────────────────────────
def draw_sidebar(slide, p_color):
    solid_rect(slide, 0, 0, SB_W, 7.5, fill=RGBColor(0xEE, 0xF2, 0xF6))
    solid_rect(slide, SB_W - 0.03, 0, 0.03, 7.5, fill=p_color)
    labels = [
        ("CAPABILITY", CAP_Y  + 0.20),
        ("STEPS",      STEP_Y + 0.20),
        ("ACTIVITIES", ACT_Y  + 1.00),
        ("ENABLERS",   EN_Y   + 0.25),
        ("NFRs",       NF_Y   + 0.18),
        ("DATA",       DT_Y   + 0.22),
    ]
    for text, y in labels:
        tb(slide, text, 0.07, y, SB_W - 0.18, 0.36,
           sz=8, bold=True, color=DIM, align=PP_ALIGN.CENTER)

# ── Flow slide factory ───────────────────────────────────────────────
def flow_slide(title, persona_label, p_color, p_dark,
               capability, trigger,
               steps, activities, enablers, nfrs, data):

    sl = prs.slides.add_slide(blank); set_bg(sl)
    draw_sidebar(sl, p_color)

    # Title bar (full header band)
    solid_rect(sl, 0, 0, 13.33, 0.46, p_color)
    tb(sl, title, CX, 0.06, 8.6, 0.36, sz=12, bold=True, color=WHITE)

    # Persona pill top-right
    round_rect(sl, 10.5, 0.07, 2.75, 0.30, fill=WHITE)
    tb(sl, persona_label, 10.5, 0.07, 2.75, 0.30,
       sz=8, bold=True, color=p_color, align=PP_ALIGN.CENTER)

    # Capability row (starts below title bar)
    solid_rect(sl, CX, CAP_Y, CW, CAP_H, fill=p_dark, border=p_color, bw=Pt(1.0))
    tb(sl, capability, CX + 0.12, CAP_Y + 0.06,
       CW - 0.24, 0.34, sz=13, bold=True, color=p_color)
    tb(sl, f"Trigger: {trigger}", CX + 0.12, CAP_Y + 0.38,
       CW - 0.24, 0.24, sz=9, italic=True, color=DIM)

    # Arrow line
    arrow_y = STEP_Y + STEP_H / 2
    solid_rect(sl, CX, arrow_y - 0.010, CW - 0.15, 0.020, fill=DIVIDER)
    tb(sl, "▶", CX + CW - 0.25, arrow_y - 0.17,
       0.28, 0.34, sz=12, color=ACCENT, align=PP_ALIGN.CENTER)

    # Step boxes
    for i, (step_name, col_x) in enumerate(zip(steps, COL_XS)):
        round_rect(sl, col_x, STEP_Y - 0.02, 0.28, 0.28, fill=p_color)
        tb(sl, str(i + 1), col_x, STEP_Y - 0.02, 0.28, 0.28,
           sz=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        shape_with_text(sl, 5, col_x + 0.04, STEP_Y, COL_W - 0.04, STEP_H,
                        fill=STEP_BG, border=p_color,
                        label_line1="STEP", label_line2=step_name,
                        sz1=7, sz2=10, c1=p_color, c2=TEXT, ml=0.30)

        for j, act_text in enumerate(activities[i]):
            ay = ACT_Y + j * (ACT_H + ACT_GAP)
            shape_with_text(sl, 5, col_x + 0.04, ay, COL_W - 0.04, ACT_H,
                            fill=ACT_BG, border=DIVIDER,
                            label_line1=None,
                            label_line2=f"→  {act_text}",
                            sz1=7, sz2=9, c2=TEXT)

    hdiv(sl, CX, EN_Y - 0.06, CW)

    # Enablers
    solid_rect(sl, CX, EN_Y, CW, EN_H, fill=BAR_BG, border=ACCENT, bw=Pt(0.75))
    tb(sl, "  |  ".join(enablers), CX + 0.12, EN_Y + 0.06,
       CW - 0.24, EN_H - 0.12, sz=10, color=TEXT)

    # NFRs
    solid_rect(sl, CX, NF_Y, CW, NF_H, fill=BAR_BG, border=DIVIDER, bw=Pt(0.75))
    tb(sl, "  |  ".join(nfrs), CX + 0.12, NF_Y + 0.06,
       CW - 0.24, NF_H - 0.12, sz=10, color=DIM)

    # Data
    solid_rect(sl, CX, DT_Y, CW, DT_H, fill=BAR_BG, border=DIVIDER, bw=Pt(0.75))
    tb(sl, "  |  ".join(data), CX + 0.12, DT_Y + 0.06,
       CW - 0.24, DT_H - 0.12, sz=10, color=DIM)

    solid_rect(sl, 0, 7.45, 13.33, 0.05, fill=p_color)


# ════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); set_bg(sl)
solid_rect(sl, 0, 0, 13.33, 0.05, ACCENT)
solid_rect(sl, 0, 2.8, 13.33, 2.2, fill=RGBColor(0xED, 0xF2, 0xF7))
tb(sl, "BIRCH WILSON", 0, 0.9, 13.33, 1.6,
   sz=72, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
tb(sl, "Assignment 3  ·  Capability Mapping & Product Backlog",
   1, 3.0, 11.33, 0.6, sz=20, color=TEXT, align=PP_ALIGN.CENTER)
tb(sl, "MET CS 634 — Agile Software Development",
   1, 3.68, 11.33, 0.44, sz=13, color=DIM, align=PP_ALIGN.CENTER)
tb(sl, "John LaCroix  ·  April 2026",
   1, 4.18, 11.33, 0.44, sz=12, color=DIM, align=PP_ALIGN.CENTER)
solid_rect(sl, 0, 7.45, 13.33, 0.05, ACCENT)

# ════════════════════════════════════════════════════════════════════
# SLIDE 2 — PRODUCT VISION
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); set_bg(sl)
solid_rect(sl, 0, 0, 13.33, 0.05, ACCENT)
tb(sl, "PRODUCT VISION", 0.5, 0.14, 12.0, 0.38, sz=9, bold=True, color=DIM)
tb(sl, ('"Connecting Fans and Prospective Partners Alike\n'
        'in Unholy Harmony — The Birch Wilson Digital Experience"'),
   0.5, 0.55, 12.33, 1.4, sz=22, bold=True, italic=True,
   color=ACCENT, align=PP_ALIGN.CENTER)
hdiv(sl, 1.2, 2.1, 10.93)

tb(sl, "SUMMARY", 0.5, 2.26, 4.0, 0.32, sz=9, bold=True, color=DIM)
tb(sl, ("Birch Wilson launched a brand new website enabling fans, fellow artists, and music "
        "producers with one common goal: opportunity. An opportunity for fans to track their "
        "new favorite rock band before they go mainstream. An opportunity for aspiring producers "
        "to find that underground act that will push them to the next level."),
   0.5, 2.60, 12.33, 1.05, sz=12, color=TEXT)

tb(sl, "PROBLEM", 0.5, 3.76, 4.0, 0.32, sz=9, bold=True, color=DIM)
tb(sl, ("Underground artists are chronically undiscovered. Bloggers miss the window. "
        "Producers can't find the talent. The Birch Wilson Digital Experience solves the "
        "discovery problem — acting as both a discovery hub and a collaboration gateway."),
   0.5, 4.10, 12.33, 0.95, sz=12, color=TEXT)

tb(sl, "SOLUTION", 0.5, 5.14, 4.0, 0.32, sz=9, bold=True, color=DIM)
tb(sl, ("The site unifies the band's Spotify, Instagram, blog, contact, and show listings "
        "into a single branded digital home — giving every visitor exactly what they came for."),
   0.5, 5.48, 12.33, 0.70, sz=12, color=TEXT)
solid_rect(sl, 0, 7.45, 13.33, 0.05, ACCENT)

# ════════════════════════════════════════════════════════════════════
# SLIDE 3 — PERSONA: DAN
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); set_bg(sl)
solid_rect(sl, 0, 0, 13.33, 0.05, DAN_C)
tb(sl, "PERSONA 01", 0.5, 0.14, 4.0, 0.32, sz=8, bold=True, color=DIM)
tb(sl, "Dan the Hardcore Blogger", 0.5, 0.46, 9.5, 0.72, sz=30, bold=True, color=TEXT)
solid_rect(sl, 0.5, 1.24, 3.5, 0.03, DAN_C)
tb(sl, "Age: 30  ·  Allston, MA  ·  Freelance Writer & Underground Music Blogger",
   0.5, 1.32, 12.0, 0.36, sz=11, color=DIM)

tb(sl, "BACKGROUND", 0.5, 1.82, 5.9, 0.32, sz=9, bold=True, color=DAN_C)
tb(sl, ("Dan works for a local underground music magazine. He bar-hops every "
        "night hunting for the next band nobody has heard of yet. He treasures "
        "the intimacy of small crowds and the purity of underground acts. He blogs "
        "about the bands he finds, building his reputation as a credible voice in "
        "the underground rock scene — always trying to get there first."),
   0.5, 2.18, 5.9, 2.1, sz=12, color=TEXT)

tb(sl, "GOALS ON THE SITE", 7.0, 1.82, 5.9, 0.32, sz=9, bold=True, color=DAN_C)
for i, g in enumerate([
    "Discover Birch Wilson before they blow up",
    "Gather raw material (bio, sound, story) to write about",
    "Contribute a review to the Birch Wilson blog",
    "Find upcoming shows to cover live",
    "Build brand as a credible underground reviewer",
]):
    tb(sl, f"→  {g}", 7.0, 2.18 + i * 0.46, 5.9, 0.42, sz=11, color=TEXT)

round_rect(sl, 0.5, 4.72, 12.33, 0.92, fill=DAN_DARK, border=DAN_C, bw=Pt(0.75))
tb(sl, ('"If there\'s a deep, unknown rock band out there, I will find it,\n'
        'and then I will tell their story!"'),
   0.75, 4.82, 11.83, 0.74, sz=14, italic=True, color=TEXT, align=PP_ALIGN.CENTER)
solid_rect(sl, 0, 7.45, 13.33, 0.05, DAN_C)

# ════════════════════════════════════════════════════════════════════
# SLIDE 4 — PERSONA: JAKE
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); set_bg(sl)
solid_rect(sl, 0, 0, 13.33, 0.05, JAKE_C)
tb(sl, "PERSONA 02", 0.5, 0.14, 4.0, 0.32, sz=8, bold=True, color=DIM)
tb(sl, "Jake the Indie Producer", 0.5, 0.46, 9.5, 0.72, sz=30, bold=True, color=TEXT)
solid_rect(sl, 0.5, 1.24, 3.5, 0.03, JAKE_C)
tb(sl, "West Springfield, MA  ·  Ghost Hit Recording  ·  Indie & Rock Specialist",
   0.5, 1.32, 12.0, 0.36, sz=11, color=DIM)

tb(sl, "BACKGROUND", 0.5, 1.82, 5.9, 0.32, sz=9, bold=True, color=JAKE_C)
tb(sl, ("Jake is a producer who specializes in finding local artists to record "
        "on a budget. A friendly introvert, he takes on indie and rock projects "
        "to build his client roster — with his eye on recording a smash hit that "
        "propels him to bigger acts. He actively scouts underground bands like "
        "Birch Wilson, looking for that next breakthrough project."),
   0.5, 2.18, 5.9, 2.1, sz=12, color=TEXT)

tb(sl, "GOALS ON THE SITE", 7.0, 1.82, 5.9, 0.32, sz=9, bold=True, color=JAKE_C)
for i, g in enumerate([
    "Assess Birch Wilson's sound and production readiness",
    "Verify the band is active and gigging (not a dead project)",
    "Understand the band's aesthetic and genre niche",
    "Reach out directly with a recording pitch",
    "Build his roster with credible underground acts",
]):
    tb(sl, f"→  {g}", 7.0, 2.18 + i * 0.46, 5.9, 0.42, sz=11, color=TEXT)

round_rect(sl, 0.5, 4.72, 12.33, 0.92, fill=JAKE_DARK, border=JAKE_C, bw=Pt(0.75))
tb(sl, '"If an indie band wants to make a record, I\'ll record them for a modest price."',
   0.75, 4.82, 11.83, 0.74, sz=14, italic=True, color=TEXT, align=PP_ALIGN.CENTER)
solid_rect(sl, 0, 7.45, 13.33, 0.05, JAKE_C)

# ════════════════════════════════════════════════════════════════════
# SLIDES 5–9 — FIVE FLOWS
# ════════════════════════════════════════════════════════════════════

# ── FLOW 1 ───────────────────────────────────────────────────────────
flow_slide(
    title         = "Flow 1  ·  Dan Discovers the Band & Gathers Blog Material",
    persona_label = "Dan — Hardcore Blogger",
    p_color = DAN_C, p_dark = DAN_DARK,
    capability    = "Discover Birch Wilson & Gather Blogging Material",
    trigger       = "Dan searches for underground rock bands to cover for his magazine",
    steps         = [
        "Find the Site",
        "Explore Band Identity",
        "Listen to Music",
        "Sign Up for Updates",
    ],
    activities    = [
        ["Search 'underground rock bands MA'",
         "Click Birch Wilson in search results"],
        ["Read band bio & origin story on About page",
         "Absorb visual aesthetic & artwork on Homepage"],
        ["Stream embedded Spotify tracks on Music page",
         "Read existing blog posts for context"],
        ["Fill out email notification form",
         "Receive welcome confirmation with Spotify link"],
    ],
    enablers      = ["SEO & search indexing", "Spotify embed integration",
                     "Email / newsletter system", "Blog CMS"],
    nfrs          = ["Discoverability", "Page load performance",
                     "Usability", "Mobile responsiveness"],
    data          = ["Band biography text", "Spotify track metadata",
                     "Genre & influences info", "Subscriber email list"],
)

# ── FLOW 2 ───────────────────────────────────────────────────────────
flow_slide(
    title         = "Flow 2  ·  Dan Pitches a Blog Contribution",
    persona_label = "Dan — Hardcore Blogger",
    p_color = DAN_C, p_dark = DAN_DARK,
    capability    = "Submit a Blog Contribution Inquiry",
    trigger       = "Dan wants to publish a review on the Birch Wilson site to build his reviewer brand",
    steps         = [
        "Review Blog Content",
        "Identify Opportunity",
        "Submit Inquiry",
        "Piece Published",
    ],
    activities    = [
        ["Navigate to Blog page",
         "Read existing posts; assess tone & style"],
        ["Find blog contribution call-to-action",
         "Review submission guidelines"],
        ["Fill out contact form (name, email, pitch idea)",
         "Submit; receive auto-confirmation email"],
        ["Band reviews & publishes Dan's piece",
         "Dan shares on his blog & social feeds"],
    ],
    enablers      = ["Blog CMS", "Contact / inquiry form",
                     "Email notification system", "Social share integration"],
    nfrs          = ["Usability", "Reliability (form submission)",
                     "Content moderation", "Accessibility"],
    data          = ["Existing blog posts", "Form fields (name, email, pitch)",
                     "Auto-reply email template", "Published article content"],
)

# ── FLOW 3 ───────────────────────────────────────────────────────────
flow_slide(
    title         = "Flow 3  ·  Dan Attends a Show & Covers It Live",
    persona_label = "Dan — Hardcore Blogger",
    p_color = DAN_C, p_dark = DAN_DARK,
    capability    = "Discover, Reserve & Attend an Upcoming Show",
    trigger       = "Dan wants to cover Birch Wilson live for his underground magazine",
    steps         = [
        "Receive Notification",
        "Review Show Details",
        "Reserve Ticket",
        "Attend & Cover",
    ],
    activities    = [
        ["Receive email alert for upcoming show",
         "Click link; land on the Shows page"],
        ["Read date, venue, and city",
         "Confirm personal availability"],
        ["Fill out ticket reservation form",
         "Submit; receive confirmation email"],
        ["Attend show; experience intimate crowd",
         "Return post-show; access band photos & bio for review"],
    ],
    enablers      = ["Email notification system", "Shows / events listing",
                     "Ticket reservation form", "Photo gallery"],
    nfrs          = ["Reliability", "Timeliness (alerts before sellout)",
                     "Usability", "Scalability"],
    data          = ["Show details (date, venue, city)", "Ticket inventory",
                     "Fan email list", "Live show photography"],
)

# ── FLOW 4 ───────────────────────────────────────────────────────────
flow_slide(
    title         = "Flow 4  ·  Jake Scouts the Band & Assesses Their Sound",
    persona_label = "Jake — Indie Producer",
    p_color = JAKE_C, p_dark = JAKE_DARK,
    capability    = "Research & Evaluate Birch Wilson as a Recording Client",
    trigger       = "Jake is searching for local indie/rock bands to record as his next project",
    steps         = [
        "Discover the Site",
        "Assess the Sound",
        "Evaluate the Brand",
        "Verify Band is Active",
    ],
    activities    = [
        ["Search 'indie rock band Massachusetts'",
         "Find & click Birch Wilson in results"],
        ["Navigate to Music page",
         "Listen to Spotify embeds; note production quality & style"],
        ["Read About page; review bio & band history",
         "Assess aesthetic fit for his studio niche"],
        ["Check Shows page for upcoming gig schedule",
         "Confirm band is actively performing (not dormant)"],
    ],
    enablers      = ["SEO & search indexing", "Spotify embed",
                     "About / bio page", "Shows / events listing"],
    nfrs          = ["Credibility (professional appearance)", "Performance",
                     "Usability", "Discoverability"],
    data          = ["Music track data", "Band biography & history",
                     "Genre & influences metadata", "Active show schedule"],
)

# ── FLOW 5 ───────────────────────────────────────────────────────────
flow_slide(
    title         = "Flow 5  ·  Jake Reaches Out to Collaborate",
    persona_label = "Jake — Indie Producer",
    p_color = JAKE_C, p_dark = JAKE_DARK,
    capability    = "Submit a Recording Collaboration Inquiry",
    trigger       = "Jake is ready to pitch a recording session to Birch Wilson",
    steps         = [
        "Navigate to Contact",
        "Complete Inquiry Form",
        "Submit & Confirm",
        "Begin Collaboration",
    ],
    activities    = [
        ["Find Contact page via site navigation",
         "Review available contact options"],
        ["Enter name, studio, email & pitch details",
         "Describe services and budget range"],
        ["Submit form; receive auto-confirmation email",
         "Inquiry logged for band review"],
        ["Band reviews inquiry; sees credible producer",
         "Band responds; recording relationship begins"],
    ],
    enablers      = ["Contact / inquiry form", "Email notification system",
                     "Form validation", "CMS for inquiry management"],
    nfrs          = ["Reliability", "Security (form data protection)",
                     "Usability", "Response time"],
    data          = ["Producer inquiry (name, studio, email, pitch)",
                     "Band contact email", "Confirmation email template",
                     "Inquiry log"],
)

# ════════════════════════════════════════════════════════════════════
# SLIDE 10 — PRODUCT ROADMAP
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); set_bg(sl)
solid_rect(sl, 0, 0, 13.33, 0.05, ACCENT)
tb(sl, "PRODUCT ROADMAP", 0.5, 0.12, 12.33, 0.56, sz=26, bold=True, color=TEXT)
round_rect(sl, 0.5, 0.76, 12.33, 0.34, fill=RGBColor(0xED, 0xF2, 0xF7),
           border=DIVIDER, bw=Pt(0.75))
tb(sl, "Trello Board:  [Insert Trello Board URL Here]",
   0.72, 0.82, 12.0, 0.26, sz=11, color=ACCENT)

cols = [
    ("Release 1\nCore Foundation",    DAN_C,  [
        "Display band biography & story",
        "Embed Spotify music player",
        "Link Instagram feed in footer",
        "Submit booking / contact form",
        "View upcoming shows listing",
        "Sign up for email notifications",
        "Display branding & artwork on Homepage",
    ]),
    ("Release 2\nFan Engagement",     ACCENT, [
        "Purchase / reserve show tickets on-site",
        "Read blog posts & band updates",
        "Browse live show photo gallery",
        "Auto-confirmation emails after submission",
        "Share pages to social platforms",
        "Welcome email with Spotify link on signup",
    ]),
    ("Release 3\nCommunity & Growth", JAKE_C, [
        "Contribute fan / press review to blog",
        "Watch embedded video (live footage)",
        "Download band press kit / EPK",
        "Browse limited merchandise store",
        "View individual band member profiles",
        "Comment on blog posts",
    ]),
]
xs = [0.3, 4.65, 9.0]; col_w = 4.0
for (col_title, col_color, items), cx in zip(cols, xs):
    round_rect(sl, cx, 1.18, col_w, 0.56, fill=col_color)
    tb(sl, col_title, cx, 1.20, col_w, 0.52,
       sz=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        iy = 1.82 + j * 0.68
        round_rect(sl, cx, iy, col_w, 0.62,
                   fill=STEP_BG, border=col_color, bw=Pt(0.75))
        tb(sl, f"→  {item}", cx + 0.1, iy + 0.08, col_w - 0.2, 0.50, sz=10, color=TEXT)
solid_rect(sl, 0, 7.45, 13.33, 0.05, ACCENT)

# ── Save ─────────────────────────────────────────────────────────────
import os
out = r"c:\Users\lacro\OneDrive\Desktop\BU\agile\LacroixJohn_Assignment3_BirchWilson.pptx"
pdf = out.replace(".pptx", ".pdf")
prs.save(out)
print(f"Saved: {out}")

# ── Export PDF via PowerPoint COM ─────────────────────────────────────
try:
    import win32com.client
    app  = win32com.client.Dispatch("PowerPoint.Application")
    deck = app.Presentations.Open(os.path.abspath(out), WithWindow=False)
    deck.SaveAs(os.path.abspath(pdf), 32)   # 32 = ppSaveAsPDF
    deck.Close()
    app.Quit()
    print(f"PDF:   {pdf}")
except Exception as e:
    print(f"PDF failed: {e}")
